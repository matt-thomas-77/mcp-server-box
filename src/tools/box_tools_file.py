import base64
import logging
from datetime import datetime
from io import BytesIO
from typing import Any

from box_ai_agents_toolkit import (
    box_file_copy,
    box_file_delete,
    box_file_info,
    box_file_lock,
    box_file_move,
    box_file_rename,
    box_file_retention_date_clear,
    box_file_retention_date_set,
    box_file_download,
    box_file_set_description,
    box_file_set_download_company,
    box_file_set_download_open,
    box_file_set_download_reset,
    box_file_tag_add,
    box_file_tag_list,
    box_file_tag_remove,
    box_file_thumbnail_download,
    box_file_thumbnail_url,
    box_file_unlock,
)
from mcp.server.fastmcp import Context
from mcp.types import ImageContent, TextContent

from tools.box_tools_generic import get_box_client

logger = logging.getLogger(__name__)

PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PPT_MIME_TYPE = "application/vnd.ms-powerpoint"
PDF_MIME_TYPE = "application/pdf"


def _extract_file_name(file_info: Any) -> str:
    """Extract file name from common Box API response shapes."""
    if not isinstance(file_info, dict):
        return ""

    top_level_name = file_info.get("name")
    if isinstance(top_level_name, str):
        return top_level_name

    nested_file = file_info.get("file")
    if isinstance(nested_file, dict):
        nested_name = nested_file.get("name")
        if isinstance(nested_name, str):
            return nested_name

    return ""


def _collect_shape_text(shape: Any) -> list[str]:
    """Collect human-readable text from a shape's text frame or table."""
    collected: list[str] = []

    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                collected.append(text)

    if getattr(shape, "has_table", False) and shape.table:
        rows: list[str] = []
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            collected.append("Table:")
            collected.extend(rows)

    return collected


# Map python-pptx content types to MIME types for images
_PPTX_IMAGE_CONTENT_TYPES: dict[str, str] = {
    "image/png": "image/png",
    "image/jpeg": "image/jpeg",
    "image/gif": "image/gif",
    "image/bmp": "image/bmp",
    "image/tiff": "image/tiff",
    "image/x-emf": "image/x-emf",
    "image/x-wmf": "image/x-wmf",
    "image/svg+xml": "image/svg+xml",
}

# Content types that LLMs can actually interpret visually
_LLM_VISIBLE_IMAGE_TYPES: set[str] = {
    "image/png",
    "image/jpeg",
    "image/gif",
    "image/webp",
}

# Image sizing constraints
_IMAGE_MAX_DIMENSION = 1024  # px – longest side after resize
_IMAGE_JPEG_QUALITY = 80
# Total base64 image budget (bytes).  The Databricks MCP proxy enforces a 4 MB
# response cap, so we reserve ~1 MB for text/metadata and allow ~3 MB for images.
_IMAGE_TOTAL_BUDGET_B64 = 3 * 1024 * 1024


def _compress_image_bytes(image_blob: bytes) -> tuple[bytes, str]:
    """Resize and compress an image to fit LLM consumption constraints.

    Returns (compressed_bytes, mime_type).
    """
    from PIL import Image as PILImage

    img = PILImage.open(BytesIO(image_blob))

    # Resize if either dimension exceeds the cap
    if max(img.size) > _IMAGE_MAX_DIMENSION:
        img.thumbnail((_IMAGE_MAX_DIMENSION, _IMAGE_MAX_DIMENSION), PILImage.LANCZOS)

    buf = BytesIO()

    # Preserve transparency as PNG; everything else becomes JPEG for size
    if img.mode in ("RGBA", "LA", "PA") or (
        img.mode == "P" and "transparency" in img.info
    ):
        img.save(buf, format="PNG", optimize=True)
        return buf.getvalue(), "image/png"

    if img.mode != "RGB":
        img = img.convert("RGB")
    img.save(buf, format="JPEG", quality=_IMAGE_JPEG_QUALITY, optimize=True)
    return buf.getvalue(), "image/jpeg"


def _collect_shape_images(shape: Any, slide_number: int) -> list[ImageContent]:
    """Extract images from a shape, returning MCP ImageContent blocks."""
    images: list[ImageContent] = []

    image_obj = getattr(shape, "image", None)
    if image_obj is None:
        return images

    try:
        content_type = image_obj.content_type
        image_blob = image_obj.blob
    except Exception:
        logger.debug(
            "Slide %d: could not read image blob from shape '%s'",
            slide_number,
            getattr(shape, "name", "unknown"),
        )
        return images

    mime_type = _PPTX_IMAGE_CONTENT_TYPES.get(content_type, content_type)

    if mime_type not in _LLM_VISIBLE_IMAGE_TYPES:
        logger.debug(
            "Slide %d: skipping non-visual image type %s",
            slide_number,
            mime_type,
        )
        return images

    try:
        compressed_blob, compressed_mime = _compress_image_bytes(image_blob)
    except Exception:
        logger.debug(
            "Slide %d: failed to compress image from shape '%s'",
            slide_number,
            getattr(shape, "name", "unknown"),
        )
        return images

    encoded = base64.b64encode(compressed_blob).decode("utf-8")
    images.append(
        ImageContent(
            type="image",
            data=encoded,
            mimeType=compressed_mime,
        )
    )

    return images


def _extract_pptx_content_from_bytes(
    file_content: bytes,
) -> dict[str, Any] | list[TextContent | ImageContent]:
    """Extract markdown text and images from a .pptx payload.

    Returns either:
    - A dict with an ``"error"`` key on failure, or
    - A list of ``TextContent`` / ``ImageContent`` blocks on success.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {
            "error": "python-pptx is not installed. Install 'python-pptx' to enable PowerPoint extraction.",
        }

    try:
        presentation = Presentation(BytesIO(file_content))
    except Exception:
        return {
            "error": "Unable to parse file as .pptx PowerPoint presentation.",
        }

    content_blocks: list[TextContent | ImageContent] = []
    markdown_parts: list[str] = []
    image_budget_remaining = _IMAGE_TOTAL_BUDGET_B64

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_markdown: list[str] = [f"## Slide {slide_number}"]

        title_shape = getattr(slide.shapes, "title", None)
        if title_shape and getattr(title_shape, "text", "").strip():
            slide_markdown.append(f"Title: {title_shape.text.strip()}")

        body_lines: list[str] = []
        slide_images: list[ImageContent] = []

        for shape in slide.shapes:
            body_lines.extend(_collect_shape_text(shape))
            if image_budget_remaining > 0:
                slide_images.extend(_collect_shape_images(shape, slide_number))

        if body_lines:
            slide_markdown.append("\n".join(f"- {line}" for line in body_lines))
        else:
            slide_markdown.append("- (No extractable slide body text found)")

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_markdown.append("Notes:")
                slide_markdown.append(notes_text)

        slide_markdown.append("")
        markdown_parts.extend(slide_markdown)

        # Emit a text block for this slide followed by its images
        content_blocks.append(
            TextContent(type="text", text="\n".join(slide_markdown).strip())
        )
        if slide_images:
            # Label each image with its slide number, respecting the budget
            for idx, img in enumerate(slide_images, start=1):
                img_b64_size = len(img.data)
                if img_b64_size > image_budget_remaining:
                    content_blocks.append(
                        TextContent(
                            type="text",
                            text=f"Slide {slide_number} — Image {idx} of {len(slide_images)} (omitted, response size limit reached)",
                        )
                    )
                    image_budget_remaining = 0
                    continue
                label = f"Slide {slide_number} — Image {idx} of {len(slide_images)}"
                content_blocks.append(TextContent(type="text", text=label))
                content_blocks.append(img)
                image_budget_remaining -= img_b64_size

    return content_blocks


def _extract_pdf_markdown_from_bytes(file_content: bytes) -> dict[str, Any]:
    """Extract markdown and metadata from a PDF payload."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return {
            "error": "pypdf is not installed. Install 'pypdf' to enable PDF extraction.",
        }

    try:
        reader = PdfReader(BytesIO(file_content))
    except Exception:
        return {
            "error": "Unable to parse file as PDF.",
        }

    markdown_parts: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        markdown_parts.append(f"## Page {page_number}")
        page_text = (page.extract_text() or "").strip()
        if page_text:
            markdown_parts.append(page_text)
        else:
            markdown_parts.append("(No extractable page text found)")
        markdown_parts.append("")

    return {
        "representation": "text/markdown",
        "page_count": len(reader.pages),
        "content": "\n".join(markdown_parts).strip(),
    }


async def box_file_info_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    Get information about a file in Box.
    Args:
        file_id (str): The ID of the file to get information about.
    return:
        dict[str, Any]: Information about the file.
    """
    box_client = get_box_client(ctx)
    return box_file_info(box_client, file_id)


async def box_file_copy_tool(
    ctx: Context,
    file_id: str,
    destination_folder_id: str,
    new_name: str | None = None,
    version_number: int | None = None,
) -> dict[str, Any]:
    """
    Copy a file to a specified destination folder in Box.
    Args:
        file_id (str): The ID of the file to copy.
        destination_folder_id (str): The ID of the destination folder.
        new_name (str, optional): Optional new name for the copied file.
        version_number (int, optional): Optional version number of the file to copy.
    Returns:
        dict[str, Any]: Dictionary containing the copied file information or error message.
    """
    box_client = get_box_client(ctx)
    return box_file_copy(
        box_client, file_id, destination_folder_id, new_name, version_number
    )


async def box_file_delete_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    Delete a file from Box.
    Args:
        file_id (str): The ID of the file to delete.
    Returns:
        dict[str, Any]: Dictionary containing success message or error.
    """
    box_client = get_box_client(ctx)
    return box_file_delete(box_client, file_id)


async def box_file_move_tool(
    ctx: Context,
    file_id: str,
    destination_folder_id: str,
) -> dict[str, Any]:
    """
    Move a file to a specified destination folder in Box.
    Args:
        file_id (str): The ID of the file to move.
        destination_folder_id (str): The ID of the destination folder.
    Returns:
        dict[str, Any]: Dictionary containing the moved file information.
    """
    box_client = get_box_client(ctx)
    return box_file_move(box_client, file_id, destination_folder_id)


async def box_file_rename_tool(
    ctx: Context,
    file_id: str,
    new_name: str,
) -> dict[str, Any]:
    """
    Rename a file in Box.
    Args:
        file_id (str): The ID of the file to rename.
        new_name (str): The new name for the file.
    Returns:
        dict[str, Any]: Dictionary containing the renamed file information.
    """
    box_client = get_box_client(ctx)
    return box_file_rename(box_client, file_id, new_name)


async def box_file_set_description_tool(
    ctx: Context,
    file_id: str,
    description: str,
) -> dict[str, Any]:
    """
    Set or update the description of a file in Box.
    Args:
        file_id (str): The ID of the file to update.
        description (str): The new description for the file.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information.
    """
    box_client = get_box_client(ctx)
    return box_file_set_description(box_client, file_id, description)


async def box_file_retention_date_set_tool(
    ctx: Context,
    file_id: str,
    retention_date: str,
) -> dict[str, Any]:
    """
    Set a retention date for a file in Box (cannot be shortened once set).
    Args:
        file_id (str): The ID of the file to update.
        retention_date (str): The retention date for the file in ISO 8601 format.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information including retention date.
    """
    box_client = get_box_client(ctx)
    # Parse the retention date string to datetime
    retention_dt = datetime.fromisoformat(retention_date.replace("Z", "+00:00"))
    return box_file_retention_date_set(box_client, file_id, retention_dt)


async def box_file_retention_date_clear_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    Clear/remove the retention date from a file in Box.
    Args:
        file_id (str): The ID of the file to update.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information.
    """
    box_client = get_box_client(ctx)
    return box_file_retention_date_clear(box_client, file_id)


async def box_file_lock_tool(
    ctx: Context,
    file_id: str,
    lock_expires_at: str | None = None,
    is_download_prevented: bool | None = None,
) -> dict[str, Any]:
    """
    Define a lock on a file to prevent it from being moved, renamed, or changed by anyone other than the lock creator.
    Args:
        file_id (str): The ID of the file to lock.
        lock_expires_at (str, optional): Optional expiration date/time for the lock in ISO 8601 format.
        is_download_prevented (bool, optional): Optional flag to prevent downloads while locked.
    Returns:
        dict[str, Any]: Dictionary containing the locked file information including lock details.
    """
    box_client = get_box_client(ctx)

    lock_expires_dt = None
    if lock_expires_at:
        lock_expires_dt = datetime.fromisoformat(lock_expires_at.replace("Z", "+00:00"))

    return box_file_lock(
        box_client, file_id, lock_expires_dt, is_download_prevented
    )


async def box_file_unlock_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    Remove a lock from a file in Box.
    Args:
        file_id (str): The ID of the file to unlock.
    Returns:
        dict[str, Any]: Dictionary containing the unlocked file information.
    """
    box_client = get_box_client(ctx)
    return box_file_unlock(box_client, file_id)


async def box_file_set_download_open_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    Allow anyone with access to the file to download it (overrides role-based download permissions).
    Args:
        file_id (str): The ID of the file to update.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information.
    """
    box_client = get_box_client(ctx)
    return box_file_set_download_open(box_client, file_id)


async def box_file_set_download_company_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    Set a file to be downloadable by company users (restricts external user downloads for viewer/editor roles).
    Args:
        file_id (str): The ID of the file to update.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information.
    """
    box_client = get_box_client(ctx)
    return box_file_set_download_company(box_client, file_id)


async def box_file_set_download_reset_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    Reset download permissions to default behavior based on collaboration roles.
    Args:
        file_id (str): The ID of the file to update.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information.
    """
    box_client = get_box_client(ctx)
    return box_file_set_download_reset(box_client, file_id)


async def box_file_tag_list_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any]:
    """
    List all tags associated with a file in Box.
    Args:
        file_id (str): The ID of the file to retrieve tags for.
    Returns:
        dict[str, Any]: Dictionary with list of tags or message if no tags found.
    """
    box_client = get_box_client(ctx)
    return box_file_tag_list(box_client, file_id)


async def box_file_tag_add_tool(
    ctx: Context,
    file_id: str,
    tag: str,
) -> dict[str, Any]:
    """
    Add a tag to a file in Box (prevents duplicates).
    Args:
        file_id (str): The ID of the file to add a tag to.
        tag (str): The tag to add.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information including tags.
    """
    box_client = get_box_client(ctx)
    return box_file_tag_add(box_client, file_id, tag)


async def box_file_tag_remove_tool(
    ctx: Context,
    file_id: str,
    tag: str,
) -> dict[str, Any]:
    """
    Remove a tag from a file in Box.
    Args:
        file_id (str): The ID of the file to remove a tag from.
        tag (str): The tag to remove.
    Returns:
        dict[str, Any]: Dictionary containing the updated file information including tags.
    """
    box_client = get_box_client(ctx)
    return box_file_tag_remove(box_client, file_id, tag)


async def box_file_thumbnail_url_tool(
    ctx: Context,
    file_id: str,
    extension: str | None = None,
    min_height: int | None = None,
    min_width: int | None = None,
    max_height: int | None = None,
    max_width: int | None = None,
) -> dict[str, Any]:
    """
    Retrieve the URL for a thumbnail image of a file.
    Args:
        file_id (str): The ID of the file.
        extension (str, optional): Image format ('png' or 'jpg', defaults to 'png').
        min_height (int, optional): Minimum height in pixels (32-320).
        min_width (int, optional): Minimum width in pixels (32-320).
        max_height (int, optional): Maximum height in pixels (32-320).
        max_width (int, optional): Maximum width in pixels (32-320).
    Returns:
        dict[str, Any]: Dictionary with thumbnail URL or message if not available.
    """
    box_client = get_box_client(ctx)
    return box_file_thumbnail_url(
        box_client, file_id, extension, min_height, min_width, max_height, max_width
    )


async def box_file_thumbnail_download_tool(
    ctx: Context,
    file_id: str,
    extension: str | None = None,
    min_height: int | None = None,
    min_width: int | None = None,
    max_height: int | None = None,
    max_width: int | None = None,
) -> dict[str, Any]:
    """
    Download the actual thumbnail image of a file.
    Args:
        file_id (str): The ID of the file.
        extension (str, optional): Image format ('png' or 'jpg', defaults to 'png').
        min_height (int, optional): Minimum height in pixels (32-320).
        min_width (int, optional): Minimum width in pixels (32-320).
        max_height (int, optional): Maximum height in pixels (32-320).
        max_width (int, optional): Maximum width in pixels (32-320).
    Returns:
        dict[str, Any]: Dictionary with thumbnail image content in base64 or error message.
    """
    box_client = get_box_client(ctx)
    result = box_file_thumbnail_download(
        box_client, file_id, extension, min_height, min_width, max_height, max_width
    )

    # If the result contains binary data, encode it as base64
    if isinstance(result, dict) and "content" in result:
        if isinstance(result["content"], bytes):
            result["content"] = base64.b64encode(result["content"]).decode("utf-8")

    return result


async def box_file_presentation_extract_tool(
    ctx: Context,
    file_id: str,
) -> dict[str, Any] | list[TextContent | ImageContent]:
    """
    Extract LLM-ready content from a PowerPoint or PDF file in Box.

    For .pptx files the response includes both markdown text **and** embedded
    slide images so that the consuming agent can visually interpret charts,
    diagrams, and other non-text content.  Each image is labelled with its
    slide number.

    For PDF files the response contains markdown text only.

    This tool only reads file bytes and does not modify the original Box file.

    Args:
        file_id (str): The ID of the file to process.

    Returns:
        Mixed content blocks (text + images) for .pptx, or a dict for PDF / errors.
    """
    box_client = get_box_client(ctx)

    file_info = box_file_info(box_client, file_id)
    file_name = _extract_file_name(file_info)

    _, file_content, mime_type = box_file_download(box_client, file_id, False, None)
    mime_type = mime_type or ""

    if mime_type == PPT_MIME_TYPE or file_name.lower().endswith(".ppt"):
        return {
            "error": "Legacy .ppt files are not supported by this extractor. Convert to .pptx and try again.",
            "mime_type": mime_type,
            "file_name": file_name,
        }

    # If metadata explicitly says this is not a supported type, fail fast.
    has_supported_mime = mime_type in {PPTX_MIME_TYPE, PPT_MIME_TYPE, PDF_MIME_TYPE}
    has_supported_filename = (
        file_name.lower().endswith(".pptx")
        or file_name.lower().endswith(".ppt")
        or file_name.lower().endswith(".pdf")
    )
    if bool(mime_type) and not has_supported_mime and not has_supported_filename:
        return {
            "error": "File is not a supported presentation format (.pptx or .pdf).",
            "mime_type": mime_type,
            "file_name": file_name,
            "file_id": file_id,
        }

    if not file_content:
        return {
            "error": "Unable to download file content from Box.",
            "mime_type": mime_type,
            "file_name": file_name,
        }

    is_pdf_hint = mime_type == PDF_MIME_TYPE or file_name.lower().endswith(".pdf")
    is_pptx_hint = mime_type == PPTX_MIME_TYPE or file_name.lower().endswith(".pptx")

    if is_pdf_hint:
        extracted = _extract_pdf_markdown_from_bytes(file_content)
        if "error" not in extracted:
            extracted["file_id"] = file_id
            extracted["file_name"] = file_name
            extracted["mime_type"] = mime_type
            return extracted
        return {
            "error": "Unable to parse file as PDF. The file may be corrupted or not a valid PDF.",
            "mime_type": mime_type,
            "file_name": file_name,
            "file_id": file_id,
        }

    # --- .pptx path: return rich content (text + images) ---
    extracted = _extract_pptx_content_from_bytes(file_content)

    # If extraction returned a dict it means an error occurred.
    if isinstance(extracted, dict):
        if extracted.get("error") in {
            "python-pptx is not installed. Install 'python-pptx' to enable PowerPoint extraction.",
        }:
            extracted["file_id"] = file_id
            extracted["file_name"] = file_name
            extracted["mime_type"] = mime_type
            return extracted

        if is_pptx_hint:
            return {
                "error": "Unable to parse file as .pptx PowerPoint presentation. The file may be corrupted or not a valid .pptx.",
                "mime_type": mime_type,
                "file_name": file_name,
                "file_id": file_id,
            }

        return {
            "error": "File is not a supported presentation format (.pptx or .pdf).",
            "mime_type": mime_type,
            "file_name": file_name,
            "file_id": file_id,
        }

    # Prepend a metadata header so the agent has file context.
    header = TextContent(
        type="text",
        text=(
            f"**File:** {file_name}\n"
            f"**File ID:** {file_id}\n"
            f"**MIME type:** {mime_type}\n"
            f"**Format:** PowerPoint (.pptx)"
        ),
    )
    return [header, *extracted]
